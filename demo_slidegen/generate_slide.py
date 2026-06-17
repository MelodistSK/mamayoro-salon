# -*- coding: utf-8 -*-
"""
例会お祝いスライド ジェネレータ（デモ用・自己完結）
  db/members.json + db/event.json（ローカルDB）を読み、指定月のお祝い対象を抽出して
  navy×gold のお祝いスライド(.pptx)を生成する。

  使い方:
    python generate_slide.py                 # event.json の月で生成 → 自動で開く
    python generate_slide.py --month 11      # 月を指定（DBから別の人が拾われる）
    python generate_slide.py -i              # 連続生成モード（月を入力するたびに生成＆オープン）
    python generate_slide.py --no-open       # 生成だけ（自動で開かない）
    python generate_slide.py --rebuild-bg    # 背景画像を作り直す

  依存: python-pptx, Pillow   （pip install python-pptx Pillow）
  ※本番システム(slide-builder)には一切依存しません。
"""
import os, sys, json, math, argparse, subprocess
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))

GOLD     = RGBColor(0xE7, 0xC4, 0x77)   # メインの金
GOLD2    = RGBColor(0xF0, 0xD4, 0x88)   # 明るい金（見出し）
GOLD_DIM = RGBColor(0x7A, 0x66, 0x3C)   # くすんだ金（罫線・仕切り）
CREAM    = RGBColor(0xF7, 0xF1, 0xE2)
DARK     = RGBColor(0x1A, 0x14, 0x08)
DIM      = RGBColor(0xB6, 0xAE, 0x9C)

FONT_MIN  = "游明朝"      # Yu Mincho
FONT_GOT  = "游ゴシック"   # Yu Gothic
FONT_MONO = "Consolas"

def load_db():
    with open(os.path.join(HERE, "db", "members.json"), encoding="utf-8") as f: members = json.load(f)
    with open(os.path.join(HERE, "db", "event.json"),   encoding="utf-8") as f: event   = json.load(f)
    return members, event

def mm(s): return int(s[:2]) if s else None
def dd(s): return int(s[3:5]) if s else None
def badge(s): return "%d/%d" % (mm(s), dd(s))

def make_bg(path, W=1920, H=1080):
    """navy 基調 ＋ 右に歯車エンブレム ＋ 上下のフェード罫飾り（モック準拠）"""
    img = Image.new("RGB", (W, H), (11, 23, 51))                # #0b1733
    # 右上から差すブライトな navy（#1d3a6e）
    glow = Image.new("L", (W, H), 0); gd = ImageDraw.Draw(glow)
    gcx, gcy, gr = int(W*0.78), int(-H*0.10), int(W*0.95)
    gd.ellipse([gcx-gr, gcy-gr, gcx+gr, gcy+gr], fill=255)
    glow = glow.filter(ImageFilter.GaussianBlur(360))
    img = Image.composite(Image.new("RGB", (W, H), (29, 58, 110)), img, glow)
    # ビネット（四隅を締める）
    vg = Image.new("L", (W, H), 0); vd = ImageDraw.Draw(vg)
    vd.ellipse([int(-W*0.16), int(-H*0.16), int(W*1.16), int(H*1.16)], fill=255)
    vg = vg.filter(ImageFilter.GaussianBlur(320))
    img = Image.composite(img, Image.new("RGB", (W, H), (5, 11, 26)), vg)

    draw = ImageDraw.Draw(img, "RGBA"); gold = (231, 196, 119)
    # 歯車エンブレム（右・半分見切れ、低不透明度）
    cx, cy, R = W*0.92, H*0.5, H*0.34; r = R*0.42; ga = gold + (24,)
    draw.ellipse([cx-R, cy-R, cx+R, cy+R], outline=ga, width=3)
    draw.ellipse([cx-r, cy-r, cx+r, cy+r], outline=ga, width=3)
    for i in range(24):
        a = i/24*2*math.pi
        draw.line([(cx+math.cos(a)*R, cy+math.sin(a)*R), (cx+math.cos(a)*(R+H*0.03), cy+math.sin(a)*(R+H*0.03))], fill=ga, width=4)
    for k in range(6):
        a = k/6*2*math.pi
        draw.line([(cx+math.cos(a)*r, cy+math.sin(a)*r), (cx+math.cos(a)*R, cy+math.sin(a)*R)], fill=ga, width=3)
    # 上下のフェード罫飾り（両端が消える金のライン）
    def hline(y, maxa=150):
        x0, x1, seg = int(W*0.06), int(W*0.94), 140
        for i in range(seg):
            t0 = i/seg; t1 = (i+1)/seg; mid = (t0+t1)/2
            xa = x0+(x1-x0)*t0; xb = x0+(x1-x0)*t1
            a = int(maxa*(1-abs(mid-0.5)*2))
            if a > 0: draw.line([(xa, y), (xb, y)], fill=gold + (a,), width=2)
    hline(int(H*0.11)); hline(int(H*0.90))
    img.save(path, "PNG")

def add_text(slide, l, t, w, h, text, font=FONT_GOT, size=14, color=CREAM,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def add_run(p, text, font, size, color, bold=False):
    r = p.add_run(); r.text = text
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return r

def rule(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.shadow.inherit = False; sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def build(month=None, rebuild_bg=False):
    members, event = load_db()
    ev_date = event["event_date"]; yr = ev_date[:4]
    M = month or int(ev_date[5:7])   # event_date は "YYYY-MM-DD"。月は 6-7 文字目から取る
    bdays = sorted([m for m in members if m["birthday"]   and mm(m["birthday"])   == M], key=lambda m: dd(m["birthday"]))
    anns  = sorted([m for m in members if m["married_on"] and mm(m["married_on"]) == M and not m["married_on_private"]], key=lambda m: dd(m["married_on"]))
    skipped = [m for m in members if m["married_on"] and mm(m["married_on"]) == M and m["married_on_private"]]

    print("=" * 54)
    print(" DB読込 : members.json (%d件) / event.json" % len(members))
    print(" 例会   : %s  %s  %s" % (event["meeting_no_kanji"], ev_date, event["venue"]))
    print(" 抽出[%d月] : 誕生日 %d名 / 結婚記念日 %d名%s"
          % (M, len(bdays), len(anns), ("  (※非公開スキップ %d)" % len(skipped)) if skipped else ""))
    for m in bdays: print("   [誕] %-5s %s" % (badge(m["birthday"]),   m["name"]))
    for m in anns:  print("   [婚] %-5s %s ・ %s" % (badge(m["married_on"]), m["name"], m["spouse_name"]))
    print("=" * 54)

    os.makedirs(os.path.join(HERE, "assets"), exist_ok=True)
    bg = os.path.join(HERE, "assets", "bg.png")
    if rebuild_bg or not os.path.exists(bg):
        print(" 背景生成 : assets/bg.png"); make_bg(bg)
    else:
        print(" 背景再利用 : assets/bg.png (--rebuild-bg で作り直し)")

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # ヘッダー
    add_text(slide, 0.92, 0.62, 11.5, 0.4, "静岡ロータリークラブ ・ %s 例会" % event["meeting_no_kanji"], FONT_GOT, 13, GOLD2)
    add_text(slide, 0.88, 1.02, 11.5, 1.2, "%d月のお祝い" % M, FONT_MIN, 52, GOLD, True)
    rule(slide, 0.92, 2.12, 2.3, 0.035, GOLD)   # 見出し下のアクセント罫

    rows = max(len(bdays), len(anns), 1)
    rowH, nameSz = (0.64, 26) if rows <= 5 else ((0.5, 19) if rows <= 7 else (0.42, 16))
    topY = 2.75

    # 2カラムの間の縦仕切り
    rule(slide, 6.62, topY + 0.05, 0.014, min(3.9, 0.72 + rows * rowH + 0.1), GOLD_DIM)

    for (lst, label, key, colx) in [(bdays, "お誕生日", "birthday", 0.95), (anns, "ご結婚記念日", "married_on", 6.98)]:
        add_text(slide, colx, topY, 4.0, 0.5, label, FONT_GOT, 18, GOLD2, True)
        rule(slide, colx + 1.95, topY + 0.21, 3.35, 0.013, GOLD_DIM)   # 見出し後ろのヘアライン
        y = topY + 0.74
        if not lst:
            add_text(slide, colx, y, 5.2, 0.5, "— 今月は該当なし —", FONT_MIN, 15, DIM)
        for m in lst:
            bh = min(0.42, rowH - 0.08)
            bd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(colx), Inches(y), Inches(0.92), Inches(bh))
            bd.shadow.inherit = False
            bd.fill.solid(); bd.fill.fore_color.rgb = GOLD; bd.line.fill.background()
            tf = bd.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
            rr = pr.add_run(); rr.text = badge(m[key]); rr.font.name = FONT_MONO; rr.font.size = Pt(13); rr.font.bold = True; rr.font.color.rgb = DARK
            # 氏名（明朝） ＋ 「さん」/配偶者（小さめ金）
            nb = slide.shapes.add_textbox(Inches(colx + 1.08), Inches(y - 0.06), Inches(4.7), Inches(rowH))
            ntf = nb.text_frame; ntf.word_wrap = True; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = 0
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.LEFT
            add_run(np, m["name"], FONT_MIN, nameSz, CREAM, True)
            if key == "birthday":
                add_run(np, "  さん", FONT_GOT, max(11, nameSz - 9), GOLD2, False)
            else:
                add_run(np, "  ・ %s さん" % m["spouse_name"], FONT_GOT, max(11, nameSz - 9), GOLD2, False)
            y += rowH

    add_text(slide, 0, 6.96, 13.333, 0.4, "%s ｜ %s" % (ev_date.replace("-", "."), event["venue"]), FONT_GOT, 11, DIM, False, PP_ALIGN.CENTER)

    os.makedirs(os.path.join(HERE, "out"), exist_ok=True)
    out = os.path.join(HERE, "out", "お祝い_%s-%02d.pptx" % (yr, M)); prs.save(out)
    print(" 生成完了 -> %s" % out)
    return out

def open_file(path):
    try:
        if os.name == "nt": os.startfile(path)                       # Windows
        elif sys.platform == "darwin": subprocess.Popen(["open", path])
        else: subprocess.Popen(["xdg-open", path])
        print(" 開いています : %s" % os.path.basename(path))
    except Exception as e:
        print(" (自動オープン不可: %s)" % e)

def interactive(no_open, rebuild_bg):
    print("\n[連続生成モード] 月(1-12)を入力して Enter。空Enter=event.jsonの月 / q=終了\n")
    first = True
    while True:
        try: s = input(" 月> ").strip()
        except (EOFError, KeyboardInterrupt): print(); break
        if s.lower() in ("q", "quit", "exit"): break
        m = int(s) if s.isdigit() and 1 <= int(s) <= 12 else None
        if s and m is None: print("  （1-12 で入力してください）"); continue
        out = build(month=m, rebuild_bg=(rebuild_bg and first)); first = False
        if not no_open: open_file(out)

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--month", type=int, default=None, help="例会の月（省略時は event.json）")
    ap.add_argument("--no-open", action="store_true", help="生成後に自動で開かない")
    ap.add_argument("-i", "--interactive", action="store_true", help="月を入力して連続生成")
    ap.add_argument("--rebuild-bg", action="store_true", help="背景画像を作り直す")
    args = ap.parse_args()
    if args.interactive:
        interactive(args.no_open, args.rebuild_bg)
    else:
        out = build(args.month, rebuild_bg=args.rebuild_bg)
        if not args.no_open: open_file(out)
